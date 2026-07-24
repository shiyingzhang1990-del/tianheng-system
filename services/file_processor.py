import os
import re
import zipfile
from flask import current_app

def process_file(file_path, file_type):
    """
    处理不同类型的文件，提取文本内容
    
    参数:
    file_path (str): 文件路径
    file_type (str): 文件类型
    
    返回:
    str: 提取的文本内容
    """
    try:
        # 如果文件类型是txt，先检查是否实际上是其他格式
        if file_type == 'txt':
            detected_type = detect_file_type(file_path)
            if detected_type != 'txt':
                file_type = detected_type
                current_app.logger.info(f"Auto-detected file type: {file_type} for {file_path}")
        
        if file_type == 'txt':
            return process_txt_file(file_path)
        elif file_type in ['doc', 'docx']:
            return process_word_file(file_path)
        elif file_type in ['ppt', 'pptx']:
            return process_ppt_file(file_path)
        elif file_type == 'pdf':
            return process_pdf_file(file_path)
        elif file_type in ['mp3', 'wav', 'ogg', 'm4a', 'flac', 'webm']:
            return process_audio_file_with_conversion(file_path, file_type)
        else:
            current_app.logger.error(f"Unsupported file type: {file_type}")
            return "Unsupported file type"
    except Exception as e:
        current_app.logger.error(f"Error processing file: {e}")
        return f"Error processing file: {str(e)}"

def process_txt_file(file_path):
    """处理TXT文件"""
    encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1', 'ascii']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
        except UnicodeDecodeError:
            continue
        except Exception as e:
            current_app.logger.error(f"Error reading file with {encoding} encoding: {e}")
            continue
    
    # 如果所有编码都失败，尝试二进制读取
    try:
        with open(file_path, 'rb') as file:
            content = file.read()
            # 尝试检测编码
            try:
                import chardet
                detected = chardet.detect(content)
                encoding = detected['encoding']
                return content.decode(encoding)
            except (ImportError, UnicodeDecodeError):
                # 如果chardet不可用或检测失败，返回部分内容
                return "文件编码无法识别，返回部分内容：" + str(content[:1000])
    except Exception as e:
        current_app.logger.error(f"Error processing file in binary mode: {e}")
        return "文件无法读取"

def process_word_file(file_path):
    """处理Word文件，无论如何都要尝试解析"""
    try:
        import docx
        import zipfile
        import xml.etree.ElementTree as ET
        
        # 检查文件扩展名
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # 首先尝试使用python-docx解析（适用于大部分情况）
        try:
            doc = docx.Document(file_path)
            full_text = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            if full_text:
                current_app.logger.info(f"Successfully processed Word file with python-docx: {file_path}")
                return '\n'.join(full_text)
            else:
                current_app.logger.warning(f"No text content found in Word file: {file_path}")
                # 即使没有文本，也尝试手动解析
                return parse_word_xml_manually(file_path)
                
        except Exception as docx_error:
            current_app.logger.warning(f"python-docx failed, trying manual XML parsing: {docx_error}")
            # 如果python-docx失败，尝试手动解析XML
            return parse_word_xml_manually(file_path)
            
    except ImportError:
        current_app.logger.error("python-docx library not installed")
        return "python-docx库未安装"
    except Exception as e:
        current_app.logger.error(f"Error processing Word file: {e}")
        # 即使出现异常，也尝试手动解析
        try:
            return parse_word_xml_manually(file_path)
        except Exception as manual_error:
            current_app.logger.error(f"Manual parsing also failed: {manual_error}")
            return f"Word文件处理失败: {str(e)}"

def process_ppt_file(file_path):
    """处理PPT文件"""
    try:
        # 这里需要安装python-pptx库
        import pptx
        
        presentation = pptx.Presentation(file_path)
        full_text = []
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        
        return '\n'.join(full_text)
    except ImportError:
        return "python-pptx library not installed"
    except Exception as e:
        current_app.logger.error(f"Error processing PPT file: {e}")
        return f"Error processing PPT file: {str(e)}"

def process_pdf_file(file_path):
    """处理PDF文件"""
    try:
        # 这里需要安装PyPDF2库
        from PyPDF2 import PdfReader
        
        reader = PdfReader(file_path)
        full_text = []
        
        for page in reader.pages:
            full_text.append(page.extract_text())
        
        return '\n'.join(full_text)
    except ImportError:
        return "PyPDF2 library not installed"
    except Exception as e:
        current_app.logger.error(f"Error processing PDF file: {e}")
        return f"Error processing PDF file: {str(e)}"

def validate_office_xml_file(file_path):
    """
    验证Office Open XML文件的有效性
    
    参数:
    file_path (str): 文件路径
    
    返回:
    bool: 文件是否有效
    """
    try:
        # 首先检查文件是否存在
        if not os.path.exists(file_path):
            current_app.logger.error(f"File does not exist: {file_path}")
            return False
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            current_app.logger.error(f"File is empty: {file_path}")
            return False
        
        # 检查文件扩展名
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in ['.docx', '.doc']:
            current_app.logger.warning(f"File extension {file_ext} may not be a Word document: {file_path}")
        
        # 检查是否为ZIP文件
        if not zipfile.is_zipfile(file_path):
            # 对于.doc文件，可能不是ZIP格式，直接返回True让后续处理
            if file_ext == '.doc':
                current_app.logger.info(f"Legacy .doc file detected, skipping ZIP validation: {file_path}")
                return True
            else:
                current_app.logger.error(f"File is not a valid ZIP file: {file_path}")
                return False
            
        with zipfile.ZipFile(file_path, 'r') as zip_file:
            file_list = zip_file.namelist()
            
            # 检查word/document.xml是否存在（这是最关键的）
            if 'word/document.xml' not in file_list:
                current_app.logger.error(f"word/document.xml missing in {file_path}")
                return False
            
            # 检查[Content_Types].xml是否存在
            if '[Content_Types].xml' not in file_list:
                current_app.logger.warning(f"[Content_Types].xml missing in {file_path}, but document.xml exists")
                # 不直接返回False，因为手动解析可能仍然有效
                return True
            
            # 验证[Content_Types].xml内容（如果存在）
            try:
                content_types = zip_file.read('[Content_Types].xml')
                import xml.etree.ElementTree as ET
                ET.fromstring(content_types)
            except Exception as e:
                current_app.logger.warning(f"Invalid [Content_Types].xml in {file_path}: {e}, but continuing with manual parsing")
                # 不直接返回False，允许手动解析
            
            return True
            
    except Exception as e:
        current_app.logger.error(f"Error validating Office XML file {file_path}: {e}")
        return False

def parse_word_xml_manually(file_path):
    """
    手动解析Word XML文件，作为python-docx的备用方案
    
    参数:
    file_path (str): 文件路径
    
    返回:
    str: 提取的文本内容
    """
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        
        # 检查文件是否为ZIP格式
        if not zipfile.is_zipfile(file_path):
            current_app.logger.warning(f"File is not a ZIP file, trying to read as text: {file_path}")
            # 如果不是ZIP文件，尝试直接读取文本
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                    # 简单提取可能的文本内容
                    import re
                    # 移除XML标签，保留文本
                    text = re.sub(r'<[^>]+>', '', content)
                    # 清理多余的空白
                    text = re.sub(r'\s+', ' ', text).strip()
                    if text:
                        return text
                    else:
                        return "无法从文档中提取文本内容"
            except Exception as text_error:
                current_app.logger.error(f"Failed to read as text: {text_error}")
                return f"文件读取失败: {str(text_error)}"
        
        with zipfile.ZipFile(file_path, 'r') as zip_file:
            file_list = zip_file.namelist()
            
            # 查找主文档文件，尝试多种可能的路径
            document_file = None
            possible_paths = [
                'word/document.xml',
                'word/document.xml.rels',
                'document.xml'
            ]
            
            for file_name in file_list:
                if any(path in file_name for path in possible_paths):
                    document_file = file_name
                    current_app.logger.info(f"Found document file: {document_file}")
                    break
            
            if not document_file:
                current_app.logger.warning(f"No standard document.xml found in {file_path}, trying alternative approach")
                # 尝试查找任何包含文本内容的XML文件
                for file_name in file_list:
                    if file_name.endswith('.xml') and ('word' in file_name or 'document' in file_name):
                        try:
                            test_xml = zip_file.read(file_name)
                            if b'<w:t>' in test_xml or b'<w:p>' in test_xml:
                                document_file = file_name
                                current_app.logger.info(f"Found alternative document file: {document_file}")
                                break
                        except:
                            continue
                
                if not document_file:
                    current_app.logger.warning(f"No suitable XML file found, trying to extract from any XML file")
                    # 如果还是找不到，尝试从任何XML文件中提取
                    for file_name in file_list:
                        if file_name.endswith('.xml'):
                            try:
                                test_xml = zip_file.read(file_name)
                                if b'<' in test_xml and b'>' in test_xml:
                                    document_file = file_name
                                    current_app.logger.info(f"Using fallback XML file: {document_file}")
                                    break
                            except:
                                continue
            
            if not document_file:
                current_app.logger.error(f"No suitable XML file found in {file_path}")
                return "未找到可解析的文档文件"
            
            # 读取主文档XML
            try:
                document_xml = zip_file.read(document_file)
                root = ET.fromstring(document_xml)
            except Exception as xml_error:
                current_app.logger.error(f"Failed to parse XML: {xml_error}")
                return f"XML解析失败: {str(xml_error)}"
            
            # 定义Word XML命名空间（多种可能的命名空间）
            namespaces = {
                'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
            }
            
            full_text = []
            
            # 方法1：尝试提取段落文本
            try:
                paragraphs = root.findall('.//w:p', namespaces)
                for para in paragraphs:
                    text_nodes = para.findall('.//w:t', namespaces)
                    para_text = ''.join([node.text or '' for node in text_nodes])
                    if para_text.strip():
                        full_text.append(para_text.strip())
            except Exception as e:
                current_app.logger.warning(f"Failed to extract paragraphs: {e}")
            
            # 方法2：如果段落提取失败，尝试查找所有文本节点
            if not full_text:
                try:
                    all_text_nodes = root.findall('.//w:t', namespaces)
                    if all_text_nodes:
                        text_content = ''.join([node.text or '' for node in all_text_nodes])
                        if text_content.strip():
                            full_text = [text_content.strip()]
                except Exception as e:
                    current_app.logger.warning(f"Failed to extract text nodes: {e}")
            
            # 方法3：如果Word命名空间失败，尝试无命名空间提取
            if not full_text:
                try:
                    # 查找所有包含文本的标签
                    for elem in root.iter():
                        if elem.text and elem.text.strip():
                            full_text.append(elem.text.strip())
                except Exception as e:
                    current_app.logger.warning(f"Failed to extract without namespaces: {e}")
            
            # 方法4：如果XML解析失败，尝试正则表达式提取
            if not full_text:
                try:
                    import re
                    xml_content = zip_file.read(document_file).decode('utf-8', errors='ignore')
                    # 提取所有文本内容
                    text_matches = re.findall(r'<w:t[^>]*>([^<]*)</w:t>', xml_content)
                    if text_matches:
                        full_text = [text.strip() for text in text_matches if text.strip()]
                except Exception as e:
                    current_app.logger.warning(f"Failed to extract with regex: {e}")
            
            if full_text:
                current_app.logger.info(f"Successfully extracted text manually from {file_path}")
                return '\n'.join(full_text)
            else:
                return "无法从文档中提取文本内容"
                
    except Exception as e:
        current_app.logger.error(f"Manual XML parsing failed for {file_path}: {e}")
        # 最后的备用方案：尝试从整个文件中提取任何可能的文本
        try:
            current_app.logger.info(f"Trying final fallback text extraction for {file_path}")
            return extract_text_from_any_format(file_path)
        except Exception as final_error:
            current_app.logger.error(f"Final fallback also failed: {final_error}")
            return f"XML解析失败: {str(e)}"

def extract_text_from_any_format(file_path):
    """
    从任何格式的文件中提取文本的最终备用方案
    
    参数:
    file_path (str): 文件路径
    
    返回:
    str: 提取的文本内容
    """
    try:
        import re
        
        # 尝试以二进制模式读取
        with open(file_path, 'rb') as f:
            content = f.read()
        
        # 尝试不同的编码
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'gbk', 'gb2312']
        text_content = None
        
        for encoding in encodings:
            try:
                text_content = content.decode(encoding, errors='ignore')
                break
            except:
                continue
        
        if not text_content:
            # 如果所有编码都失败，使用latin-1作为最后手段
            text_content = content.decode('latin-1', errors='ignore')
        
        # 提取所有可能的文本内容
        extracted_texts = []
        
        # 方法1：提取XML标签中的文本
        xml_texts = re.findall(r'<[^>]*>([^<]+)</[^>]*>', text_content)
        if xml_texts:
            extracted_texts.extend([text.strip() for text in xml_texts if text.strip()])
        
        # 方法2：提取引号中的文本
        quoted_texts = re.findall(r'"([^"]+)"', text_content)
        if quoted_texts:
            extracted_texts.extend([text.strip() for text in quoted_texts if text.strip()])
        
        # 方法3：提取看起来像中文或英文的文本
        chinese_texts = re.findall(r'[\u4e00-\u9fff]+', text_content)
        if chinese_texts:
            extracted_texts.extend([text.strip() for text in chinese_texts if text.strip()])
        
        english_texts = re.findall(r'[a-zA-Z]{3,}', text_content)
        if english_texts:
            extracted_texts.extend([text.strip() for text in english_texts if text.strip()])
        
        # 方法4：提取数字和符号组合
        number_texts = re.findall(r'[0-9]+[a-zA-Z\u4e00-\u9fff]*', text_content)
        if number_texts:
            extracted_texts.extend([text.strip() for text in number_texts if text.strip()])
        
        # 去重并过滤
        unique_texts = []
        seen = set()
        for text in extracted_texts:
            if text not in seen and len(text) > 1:
                unique_texts.append(text)
                seen.add(text)
        
        if unique_texts:
            current_app.logger.info(f"Successfully extracted {len(unique_texts)} text fragments using fallback method")
            return '\n'.join(unique_texts)
        else:
            return "无法从文档中提取任何文本内容"
            
    except Exception as e:
        current_app.logger.error(f"Final fallback text extraction failed: {e}")
        return f"最终文本提取失败: {str(e)}"

def detect_file_type(file_path):
    """
    基于文件内容检测文件类型
    
    参数:
    file_path (str): 文件路径
    
    返回:
    str: 检测到的文件类型
    """
    try:
        # 检查是否是ZIP文件（DOCX和PPTX都是ZIP格式）
        if zipfile.is_zipfile(file_path):
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                file_list = zip_file.namelist()
                
                # 检查是否是DOCX文件
                if 'word/document.xml' in file_list:
                    return 'docx'
                
                # 检查是否是PPTX文件
                if any(f.startswith('ppt/slides/') for f in file_list):
                    return 'pptx'
        
        # 检查是否是PDF文件
        with open(file_path, 'rb') as f:
            header = f.read(8)
            if header.startswith(b'%PDF-'):
                return 'pdf'
        
        # 默认返回txt
        return 'txt'
        
    except Exception as e:
        current_app.logger.warning(f"Error detecting file type for {file_path}: {e}")
        return 'txt'

def process_audio_file_with_conversion(file_path, file_type):
    """
    处理音频文件上传：转换为PCM并保存，然后进行语音识别
    
    参数:
    file_path (str): 音频文件路径
    file_type (str): 音频文件类型
    
    返回:
    str: 识别出的文本内容
    """
    try:
        current_app.logger.info(f"Processing audio upload: {file_path}, type: {file_type}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        if file_size > 10 * 1024 * 1024:  # 10MB
            current_app.logger.warning(f"Audio file too large: {file_size} bytes")
            return f"音频文件过大（{file_size/1024/1024:.1f}MB），百度语音识别API限制文件大小不超过10MB。请压缩音频文件后重试。"
        
        if file_size < 1000:  # 小于1KB
            current_app.logger.warning(f"Audio file too small: {file_size} bytes")
            return f"音频文件过小（{file_size}字节），可能是空文件或损坏文件。请检查音频文件完整性。"
        
        # 步骤1: 转换音频为PCM格式并保存
        pcm_file_path = convert_and_save_pcm(file_path, file_type)
        if not pcm_file_path:
            # 转换失败，返回格式特定的建议
            return get_conversion_failure_message(file_path, file_type)
        
        current_app.logger.info(f"Audio converted to PCM: {pcm_file_path}")
        
        # 步骤2: 使用PCM文件进行语音识别
        text_result = call_baidu_speech_recognition_pcm(pcm_file_path)
        
        if text_result and text_result.strip():
            current_app.logger.info(f"Successfully transcribed audio file: {file_path}")
            current_app.logger.info(f"Transcription result length: {len(text_result)} characters")
            
            # 添加音频文件标识
            formatted_result = f"=== 语音转录内容 ===\n{text_result.strip()}\n=== 转录结束 ==="
            return formatted_result
        else:
            # 语音识别失败或返回空结果
            current_app.logger.warning(f"Speech recognition failed or returned empty result for: {file_path}")
            
            # 提供详细的错误信息和建议
            return get_empty_recognition_message(file_path, file_size)
    
    except Exception as e:
        current_app.logger.error(f"Error processing audio file {file_path}: {e}")
        return f"音频文件处理出错：{str(e)}\n\n建议：请检查音频文件格式是否正确，或尝试转换为标准WAV格式后重新上传。"

def convert_and_save_pcm(original_file_path, file_type):
    """
    将音频文件转换为PCM格式并保存
    
    参数:
    original_file_path (str): 原始音频文件路径
    file_type (str): 音频文件类型
    
    返回:
    str: PCM文件路径，失败返回None
    """
    try:
        import os
        
        # 生成PCM文件路径（在同一目录下，添加.pcm扩展名）
        base_name = os.path.splitext(original_file_path)[0]
        pcm_file_path = f"{base_name}.pcm"
        
        # 优先处理WAV格式
        if file_type.lower() == 'wav':
            current_app.logger.info("Converting WAV to PCM using pure Python...")
            pcm_data = wav_to_pcm_data(original_file_path)
            
            if pcm_data:
                # 保存PCM数据到文件
                with open(pcm_file_path, 'wb') as pcm_file:
                    pcm_file.write(pcm_data)
                current_app.logger.info(f"PCM file saved: {pcm_file_path}, size: {len(pcm_data)} bytes")
                return pcm_file_path
            else:
                current_app.logger.warning("WAV to PCM conversion failed")
                return None
        
        # 处理其他格式，尝试pydub转换
        else:
            current_app.logger.info(f"Converting {file_type} to PCM using pydub...")
            
            try:
                from pydub import AudioSegment
                current_app.logger.info(f"pydub库已加载，开始转换{file_type}格式音频文件")
                
                # 加载音频文件 - 使用pydub + FFmpeg支持所有主流格式
                if file_type.lower() == 'wav':
                    current_app.logger.info("使用pydub.from_wav()加载WAV文件")
                    audio = AudioSegment.from_wav(original_file_path)
                elif file_type.lower() == 'mp3':
                    current_app.logger.info("使用pydub.from_mp3()加载MP3文件（通过FFmpeg）")
                    audio = AudioSegment.from_mp3(original_file_path)
                elif file_type.lower() == 'm4a':
                    current_app.logger.info("使用pydub.from_file()加载M4A文件（通过FFmpeg）")
                    audio = AudioSegment.from_file(original_file_path, format="m4a")
                elif file_type.lower() == 'ogg':
                    current_app.logger.info("使用pydub.from_ogg()加载OGG文件（通过FFmpeg）")
                    audio = AudioSegment.from_ogg(original_file_path)
                elif file_type.lower() == 'webm':
                    current_app.logger.info("使用pydub.from_file()加载WebM文件（通过FFmpeg）")
                    audio = AudioSegment.from_file(original_file_path, format="webm")
                elif file_type.lower() == 'flac':
                    current_app.logger.info("使用pydub.from_file()加载FLAC文件（通过FFmpeg）")
                    audio = AudioSegment.from_file(original_file_path, format="flac")
                else:
                    current_app.logger.info(f"使用pydub.from_file()加载{file_type}文件（通过FFmpeg）")
                    audio = AudioSegment.from_file(original_file_path)
                
                current_app.logger.info(f"原始音频: {audio.frame_rate}Hz, {audio.channels}声道, {audio.sample_width*8}位, {len(audio)/1000:.1f}秒")
                
                # 检查音频时长
                duration_seconds = len(audio) / 1000.0
                if duration_seconds < 1.0:
                    current_app.logger.warning(f"Audio too short: {duration_seconds} seconds")
                    return None
                elif duration_seconds > 60.0:
                    current_app.logger.warning(f"Audio too long: {duration_seconds} seconds, truncating to 60 seconds")
                    audio = audio[:60000]
                
                # 转换为百度API要求的格式：16kHz, 单声道, 16位
                current_app.logger.info("转换音频格式为百度API要求: 16kHz, 单声道, 16位")
                audio = audio.set_frame_rate(16000)
                audio = audio.set_channels(1)
                audio = audio.set_sample_width(2)
                
                current_app.logger.info(f"转换后音频: 16000Hz, 1声道, 16位, {len(audio)/1000:.1f}秒")
                
                # 导出为原始PCM数据
                pcm_data = audio.raw_data
                current_app.logger.info(f"提取PCM数据: {len(pcm_data)} bytes")
                
                # 保存PCM数据到文件
                with open(pcm_file_path, 'wb') as pcm_file:
                    pcm_file.write(pcm_data)
                
                current_app.logger.info(f"✅ pydub转换完成，PCM文件已保存: {pcm_file_path}, size: {len(pcm_data)} bytes")
                return pcm_file_path
                
            except ImportError as import_error:
                current_app.logger.error(f"pydub库未安装: {import_error}")
                current_app.logger.error("请运行: pip install pydub")
                return None
            except Exception as pydub_error:
                current_app.logger.error(f"pydub conversion failed: {pydub_error}")
                # 检查是否是FFmpeg相关问题
                if 'ffmpeg' in str(pydub_error).lower() or 'ffprobe' in str(pydub_error).lower():
                    current_app.logger.error("FFmpeg相关错误，可能是路径配置问题")
                    current_app.logger.error("请确保FFmpeg已正确安装并添加到系统PATH")
                else:
                    current_app.logger.error("音频文件可能损坏或格式异常")
                return None
    
    except Exception as e:
        current_app.logger.error(f"Error converting audio to PCM: {e}")
        return None

def call_baidu_speech_recognition_pcm(pcm_file_path):
    """
    使用PCM文件调用百度语音识别API
    
    参数:
    pcm_file_path (str): PCM文件路径
    
    返回:
    str: 识别结果文本，失败返回None
    """
    try:
        from aip import AipSpeech
        
        # 获取配置
        app_id = current_app.config.get('BAIDU_APP_ID')
        api_key = current_app.config.get('BAIDU_API_KEY')
        secret_key = current_app.config.get('BAIDU_SECRET_KEY')
        
        if not all([app_id, api_key, secret_key]):
            current_app.logger.error("Baidu Speech API configuration missing")
            return None
        
        # 初始化AipSpeech对象
        client = AipSpeech(app_id, api_key, secret_key)
        
        # 读取PCM文件
        with open(pcm_file_path, 'rb') as pcm_file:
            pcm_data = pcm_file.read()
        
        current_app.logger.info(f"Calling Baidu ASR API with PCM file: {pcm_file_path}, size: {len(pcm_data)} bytes")
        
        # 调用百度语音识别API
        result = client.asr(
            pcm_data,
            'pcm',    # 使用PCM格式
            16000,    # 16kHz采样率
            {
                'dev_pid': 1537,  # 普通话(支持简单的英文识别)
            }
        )
        
        current_app.logger.info(f"Baidu API response: {result}")
        
        # 解析结果
        if result.get('err_no') == 0:
            text_result = ''.join(result.get('result', []))
            current_app.logger.info(f"Speech recognition successful: {text_result[:100]}...")
            return text_result if text_result.strip() else None
        else:
            error_msg = result.get('err_msg', 'Unknown error')
            error_code = result.get('err_no', 'Unknown code')
            current_app.logger.error(f"Baidu Speech Recognition error [{error_code}]: {error_msg}")
            return None
            
    except ImportError:
        current_app.logger.error("baidu-aip library not installed")
        return None
    except Exception as e:
        current_app.logger.error(f"Error calling Baidu Speech API with PCM: {e}")
        return None

def get_conversion_failure_message(file_path, file_type):
    """获取转换失败的消息"""
    try:
        file_info = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        
        # 根据文件类型提供特定的建议（FFmpeg已安装，支持所有格式）
        format_specific_messages = {
            'wav': {
                'description': 'WAV是标准的无损音频格式',
                'requirements': 'pydub库（已安装）',
                'solution': 'pydub原生支持，最稳定的格式，推荐使用。'
            },
            'mp3': {
                'description': 'MP3是最常见的音频格式',
                'requirements': 'pydub库 + FFmpeg工具（已安装）',
                'solution': '通过FFmpeg支持，应该可以正常处理。'
            },
            'm4a': {
                'description': 'M4A是Apple设备常用的音频格式',
                'requirements': 'pydub库 + FFmpeg工具（已安装）',
                'solution': '通过FFmpeg支持，适合来自iPhone/iPad的录音。'
            },
            'webm': {
                'description': 'WebM是现代浏览器录音的格式',
                'requirements': 'pydub库 + FFmpeg工具（已安装）',
                'solution': '通过FFmpeg支持，适合浏览器录音文件。'
            },
            'ogg': {
                'description': 'OGG是开源的音频格式',
                'requirements': 'pydub库 + FFmpeg工具（已安装）',
                'solution': '通过FFmpeg支持，开源格式。'
            },
            'flac': {
                'description': 'FLAC是无损音频格式',
                'requirements': 'pydub库 + FFmpeg工具（已安装）',
                'solution': '通过FFmpeg支持，高质量无损格式。'
            }
        }
        
        format_info = format_specific_messages.get(file_type.lower(), {
            'description': f'{file_type.upper()}格式',
            'requirements': 'pydub库 + FFmpeg工具（已安装）',
            'solution': f'通过FFmpeg支持，理论上应该可以处理{file_type.upper()}格式'
        })
        
        return f"""音频文件 '{file_info}' ({file_type.upper()}格式, {file_size/1024:.1f}KB) 上传成功，但pydub转换失败。

📋 格式信息：
• {format_info['description']}
• 处理要求：{format_info['requirements']}

❌ 转换失败可能原因：
• 音频文件损坏或格式异常
• FFmpeg路径配置问题（虽然已安装）
• 文件编码不被支持
• 音频时长过短或过长

🔧 解决方案：
1. 🎯 首选方案：
   • 使用录音功能直接录制WAV格式（最稳定）
   • 将文件转换为WAV格式后重新上传
   
2. 🛠️ 技术排查：
   • 检查音频文件是否完整无损坏
   • 尝试使用其他工具播放该文件
   • 确认文件不是空文件或过短
   
3. 🔄 格式转换：
   • 使用在线工具转换为WAV：CloudConvert、Convertio等
   • 直接手动输入音频内容

 技术提示：{format_info['solution']}

✅ 推荐格式：WAV (16kHz, 单声道, 16位) - 最佳兼容性"""
    except:
        return f"音频文件({file_type})格式转换失败。FFmpeg已安装，但可能存在配置问题，建议转换为WAV格式以获得最佳兼容性。"

def get_empty_recognition_message(file_path, file_size):
    """获取空识别结果的消息"""
    try:
        file_info = os.path.basename(file_path)
        return f"""音频文件 '{file_info}' (大小: {file_size/1024:.1f}KB) 已成功上传并处理。

🔍 语音识别结果：
• 百度API响应成功，但未识别到可理解的语音内容
• 可能的原因：录音为静音、音量过小、背景噪音过大或语音不清晰

 建议改进：
1. 🎤 重新录音：确保在安静环境中清晰发音
2. 🔊 检查音量：确保录音音量适中，避免过小或过大
3. ⏱️ 录音时长：建议录音时长在3-30秒之间
4. 📝 手动输入：直接在文本框中输入要分析的内容

🛠️ 录音技巧：
• 距离麦克风15-20厘米
• 语速适中，发音清晰
• 避免背景噪音和回声"""
    except:
        return f"音频文件已上传，但语音识别未返回有效内容。建议重新录音或手动输入文本内容。"

def process_audio_file(file_path, file_type):
    """
    处理音频文件，使用百度语音识别API转换为文本
    
    参数:
    file_path (str): 音频文件路径
    file_type (str): 音频文件类型
    
    返回:
    str: 识别出的文本内容
    """
    try:
        current_app.logger.info(f"Processing audio file: {file_path}, type: {file_type}")
        
        # 检查文件大小（百度API要求音频文件不超过10MB）
        file_size = os.path.getsize(file_path)
        if file_size > 10 * 1024 * 1024:  # 10MB
            current_app.logger.warning(f"Audio file too large: {file_size} bytes")
            return f"音频文件过大（{file_size/1024/1024:.1f}MB），百度语音识别API限制文件大小不超过10MB。请压缩音频文件后重试。"
        
        # 检查文件时长（估算，避免过长的音频）
        if file_size < 1000:  # 小于1KB，可能是空文件或损坏文件
            current_app.logger.warning(f"Audio file too small: {file_size} bytes")
            return f"音频文件过小（{file_size}字节），可能是空文件或损坏文件。请检查音频文件完整性。"
        
        # 尝试使用百度语音识别
        text_result = call_baidu_speech_recognition(file_path, file_type)
        
        if text_result and text_result.strip():
            current_app.logger.info(f"Successfully transcribed audio file: {file_path}")
            current_app.logger.info(f"Transcription result length: {len(text_result)} characters")
            
            # 添加音频文件标识
            formatted_result = f"=== 语音转录内容 ===\n{text_result.strip()}\n=== 转录结束 ==="
            return formatted_result
        else:
            # 如果语音识别失败，返回详细的提示信息
            current_app.logger.warning(f"Speech recognition failed or returned empty result for: {file_path}")
            
            # 提供更详细的错误信息和建议
            return f"""音频文件已上传（{file_type}格式），但语音识别暂时无法处理此文件。

可能的原因：
1. 音频格式不兼容（推荐使用WAV格式）
2. 音频质量较差或背景噪音过大
3. 录音时长过短或过长
4. 网络连接问题

建议解决方案：
1. 转换为WAV格式（16kHz采样率，单声道）
2. 在安静环境中重新录音
3. 稍后重试或手动转录内容

文件路径：{file_path}
文件大小：{file_size/1024:.1f}KB"""
    
    except Exception as e:
        current_app.logger.error(f"Error processing audio file {file_path}: {e}")
        return f"音频文件处理出错：{str(e)}\n\n建议：请检查音频文件格式是否正确，或尝试转换为标准WAV格式后重新上传。"

def wav_to_pcm_data(wav_file_path):
    """
    将WAV文件转换为PCM数据（纯Python实现，无需FFmpeg）
    
    参数:
    wav_file_path (str): WAV文件路径
    
    返回:
    bytes: PCM音频数据，如果失败返回None
    """
    try:
        import wave
        # 方法1：使用wave库读取WAV文件
        with wave.open(wav_file_path, 'rb') as wav_file:
            # 获取音频参数
            channels = wav_file.getnchannels()
            sample_width = wav_file.getsampwidth() 
            frame_rate = wav_file.getframerate()
            frames = wav_file.getnframes()
            
            current_app.logger.info(f"WAV文件信息: {channels}声道, {sample_width*8}位, {frame_rate}Hz, {frames}帧")
            
            # 计算音频时长
            duration = frames / frame_rate if frame_rate > 0 else 0
            current_app.logger.info(f"音频时长: {duration:.2f}秒")
            
            # 检查音频时长
            if duration < 0.5:
                current_app.logger.warning(f"音频时长过短: {duration:.2f}秒，可能影响识别效果")
            elif duration > 60:
                current_app.logger.warning(f"音频时长过长: {duration:.2f}秒，将截取前60秒")
                # 截取前60秒
                max_frames = int(60 * frame_rate)
                frames = min(frames, max_frames)
            
            # 读取PCM数据
            pcm_data = wav_file.readframes(frames)
            
            # 检查音频是否为静音
            if check_audio_silence(pcm_data, sample_width):
                current_app.logger.warning("检测到音频可能为静音或音量过低")
            
            # 百度API要求16kHz单声道16位
            if frame_rate != 16000 or channels != 1 or sample_width != 2:
                current_app.logger.info(f"需要转换音频格式: {frame_rate}Hz -> 16000Hz, {channels}ch -> 1ch, {sample_width*8}bit -> 16bit")
                pcm_data = convert_audio_format(pcm_data, channels, sample_width, frame_rate)
            
            current_app.logger.info(f"PCM数据大小: {len(pcm_data)} bytes")
            return pcm_data
            
    except Exception as e:
        current_app.logger.error(f"WAV to PCM conversion failed: {e}")
        
        # 方法2：直接读取WAV文件并去除头部（备用方案）
        try:
            current_app.logger.info("尝试直接读取WAV文件方式...")
            with open(wav_file_path, 'rb') as f:
                # 跳过44字节的WAV文件头
                f.seek(44)
                pcm_data = f.read()
                
            current_app.logger.info(f"直接读取PCM数据大小: {len(pcm_data)} bytes")
            
            # 简单的静音检查
            if len(pcm_data) > 1000:
                # 检查前1000字节是否全为0或接近0
                sample_data = pcm_data[:1000]
                max_val = max(abs(int.from_bytes(sample_data[i:i+2], byteorder='little', signed=True)) 
                             for i in range(0, len(sample_data)-1, 2))
                if max_val < 100:  # 音量很低
                    current_app.logger.warning("音频音量可能过低，建议重新录音")
            
            return pcm_data
            
        except Exception as e2:
            current_app.logger.error(f"Direct WAV reading failed: {e2}")
            return None

def check_audio_silence(pcm_data, sample_width):
    """
    检查音频是否为静音或音量过低
    
    参数:
    pcm_data (bytes): PCM音频数据
    sample_width (int): 采样宽度（字节）
    
    返回:
    bool: True表示可能为静音
    """
    try:
        import numpy as np
        
        if len(pcm_data) < 100:
            return True
        
        # 根据采样宽度解析音频数据
        if sample_width == 1:
            audio_array = np.frombuffer(pcm_data, dtype=np.uint8)
            # 转换为有符号数据
            audio_array = audio_array.astype(np.int16) - 128
        elif sample_width == 2:
            audio_array = np.frombuffer(pcm_data, dtype=np.int16)
        elif sample_width == 4:
            audio_array = np.frombuffer(pcm_data, dtype=np.int32)
        else:
            return False
        
        # 计算音频的RMS（均方根）值
        rms = np.sqrt(np.mean(audio_array.astype(np.float64) ** 2))
        max_amplitude = np.max(np.abs(audio_array))
        
        current_app.logger.info(f"音频分析: RMS={rms:.2f}, Max={max_amplitude}")
        
        # 设置阈值判断是否为静音
        silence_threshold = 100 if sample_width == 2 else 50
        
        return rms < silence_threshold and max_amplitude < silence_threshold * 2
        
    except Exception as e:
        current_app.logger.warning(f"Audio silence check failed: {e}")
        return False

def convert_audio_format(pcm_data, channels, sample_width, frame_rate):
    """
    转换音频格式为百度API要求的16kHz单声道16位
    
    参数:
    pcm_data (bytes): 原始PCM数据
    channels (int): 声道数
    sample_width (int): 采样宽度（字节）
    frame_rate (int): 采样率
    
    返回:
    bytes: 转换后的PCM数据
    """
    try:
        import numpy as np
        
        # 将字节数据转换为numpy数组
        if sample_width == 1:
            dtype = np.uint8
        elif sample_width == 2:
            dtype = np.int16
        elif sample_width == 4:
            dtype = np.int32
        else:
            current_app.logger.warning(f"Unsupported sample width: {sample_width}")
            return pcm_data
        
        # 解码PCM数据
        audio_data = np.frombuffer(pcm_data, dtype=dtype)
        
        # 处理多声道（转为单声道）
        if channels > 1:
            audio_data = audio_data.reshape(-1, channels)
            # 取平均值转为单声道
            audio_data = np.mean(audio_data, axis=1).astype(dtype)
            current_app.logger.info(f"转换为单声道: {channels} -> 1")
        
        # 重采样到16kHz（简单的降采样）
        if frame_rate != 16000:
            # 简单的重采样：按比例取样
            ratio = frame_rate / 16000
            if ratio > 1:  # 降采样
                indices = np.arange(0, len(audio_data), ratio).astype(int)
                audio_data = audio_data[indices]
                current_app.logger.info(f"降采样: {frame_rate}Hz -> 16000Hz")
            else:  # 升采样（重复样本）
                repeat_count = int(1 / ratio)
                audio_data = np.repeat(audio_data, repeat_count)
                current_app.logger.info(f"升采样: {frame_rate}Hz -> 16000Hz")
        
        # 确保数据类型为int16
        if dtype != np.int16:
            # 规范化到int16范围
            if dtype == np.uint8:
                audio_data = ((audio_data.astype(np.float32) - 128) / 128 * 32767).astype(np.int16)
            elif dtype == np.int32:
                audio_data = (audio_data / 65536).astype(np.int16)
            current_app.logger.info(f"转换数据类型: {dtype} -> int16")
        
        # 转换回字节
        converted_data = audio_data.tobytes()
        current_app.logger.info(f"格式转换完成: {len(converted_data)} bytes")
        
        return converted_data
        
    except Exception as e:
        current_app.logger.error(f"Audio format conversion failed: {e}")
        return pcm_data  # 返回原始数据

def pcm_to_wav_data(pcm_data, channels=1, sample_width=2, frame_rate=16000):
    """
    将PCM数据转换为WAV格式（纯Python实现）
    
    参数:
    pcm_data (bytes): PCM音频数据
    channels (int): 声道数
    sample_width (int): 采样宽度（字节）
    frame_rate (int): 采样率
    
    返回:
    bytes: WAV格式的音频数据
    """
    try:
        import io
        
        # 创建内存中的WAV文件
        wav_buffer = io.BytesIO()
        
        import wave
        with wave.open(wav_buffer, 'wb') as wav_file:
            wav_file.setnchannels(channels)
            wav_file.setsampwidth(sample_width)
            wav_file.setframerate(frame_rate)
            wav_file.writeframes(pcm_data)
        
        wav_data = wav_buffer.getvalue()
        current_app.logger.info(f"PCM转WAV完成: {len(wav_data)} bytes")
        
        return wav_data
        
    except Exception as e:
        current_app.logger.error(f"PCM to WAV conversion failed: {e}")
        return None

def call_baidu_speech_recognition(file_path, file_type):
    """
    调用百度语音识别API，支持音频格式转换
    
    参数:
    file_path (str): 音频文件路径
    file_type (str): 音频文件类型
    
    返回:
    str: 识别结果文本，失败返回None
    """
    try:
        from aip import AipSpeech
        from pydub import AudioSegment
        import tempfile
        import os
        import numpy as np
        import wave
        import struct
        
        # 获取配置
        app_id = current_app.config.get('BAIDU_APP_ID')
        api_key = current_app.config.get('BAIDU_API_KEY')
        secret_key = current_app.config.get('BAIDU_SECRET_KEY')
        
        if not all([app_id, api_key, secret_key]):
            current_app.logger.error("Baidu Speech API configuration missing")
            return None
        
        # 初始化AipSpeech对象
        client = AipSpeech(app_id, api_key, secret_key)
        
        # 智能音频格式处理 - 优先使用纯Python方案
        converted_file_path = None
        audio_data = None
        
        try:
            current_app.logger.info(f"Processing audio file: {file_path}, format: {file_type}")
            
            # 优先处理WAV格式 - 使用纯Python方案
            if file_type.lower() == 'wav':
                current_app.logger.info("使用纯Python WAV处理方案...")
                audio_data = wav_to_pcm_data(file_path)
                
                if audio_data:
                    current_app.logger.info("✅ 纯Python WAV转换成功")
                else:
                    current_app.logger.warning("❌ 纯Python WAV转换失败，尝试pydub方案")
            
            # 如果WAV处理失败或非WAV格式，尝试pydub方案
            if not audio_data:
                current_app.logger.info("尝试pydub音频转换方案...")
                
                try:
                    if file_type.lower() == 'wav':
                        audio = AudioSegment.from_wav(file_path)
                        current_app.logger.info(f"Original WAV: {audio.frame_rate}Hz, {audio.channels}ch, {audio.sample_width*8}bit")
                    elif file_type.lower() == 'mp3':
                        audio = AudioSegment.from_mp3(file_path)
                    elif file_type.lower() == 'm4a':
                        audio = AudioSegment.from_file(file_path, format="m4a")
                    elif file_type.lower() == 'ogg':
                        audio = AudioSegment.from_ogg(file_path)
                    elif file_type.lower() == 'webm':
                        current_app.logger.info("Processing WebM format - converting to WAV")
                        audio = AudioSegment.from_file(file_path, format="webm")
                    elif file_type.lower() == 'flac':
                        audio = AudioSegment.from_file(file_path, format="flac")
                    else:
                        audio = AudioSegment.from_file(file_path)
                    
                    # 检查音频时长
                    duration_seconds = len(audio) / 1000.0
                    current_app.logger.info(f"Audio duration: {duration_seconds:.2f} seconds")
                    
                    if duration_seconds < 1.0:
                        current_app.logger.warning(f"Audio too short: {duration_seconds} seconds")
                        return None
                    elif duration_seconds > 60.0:
                        current_app.logger.warning(f"Audio too long: {duration_seconds} seconds, truncating to 60 seconds")
                        audio = audio[:60000]
                    
                    # 转换为百度API要求的格式
                    audio = audio.set_frame_rate(16000)
                    audio = audio.set_channels(1)
                    audio = audio.set_sample_width(2)
                    
                    current_app.logger.info(f"Converted: 16000Hz, 1ch, 16bit, {len(audio)/1000:.2f}s")
                    
                    # 创建临时文件
                    with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as temp_file:
                        converted_file_path = temp_file.name
                        
                    audio.export(converted_file_path, format="wav")
                    current_app.logger.info(f"Audio converted to PCM WAV: {converted_file_path}")
                    
                    # 读取转换后的音频数据
                    with open(converted_file_path, 'rb') as fp:
                        audio_data = fp.read()
                    
                    current_app.logger.info(f"✅ pydub转换成功: {len(audio_data)} bytes")
                    
                except Exception as pydub_error:
                    current_app.logger.warning(f"❌ pydub转换失败: {pydub_error}")
                    
                    # 最后的备用方案：直接读取WAV文件
                    if file_type.lower() == 'wav':
                        current_app.logger.info("尝试直接读取WAV文件...")
                        try:
                            with open(file_path, 'rb') as fp:
                                audio_data = fp.read()
                            
                            # 基本格式检查
                            if len(audio_data) < 44:
                                current_app.logger.error("WAV file too short")
                                return None
                            
                            if not audio_data.startswith(b'RIFF') or b'WAVE' not in audio_data[:12]:
                                current_app.logger.error("Invalid WAV file format")
                                return None
                            
                            current_app.logger.info(f"✅ 直接读取WAV成功: {len(audio_data)} bytes")
                            
                        except Exception as direct_error:
                            current_app.logger.error(f"❌ 直接读取WAV失败: {direct_error}")
                            return None
                    
            # 检查是否成功获取音频数据
            if not audio_data:
                # 提供格式特定的建议
                current_app.logger.error(f"All conversion methods failed for {file_type} format")
                
                # 为不同格式提供特定的处理建议
                try:
                    file_info = os.path.basename(file_path)
                    file_size = os.path.getsize(file_path)
                    
                    format_messages = {
                        'webm': f"""WebM音频文件 '{file_info}' (大小: {file_size/1024:.1f}KB) 已上传成功。

📋 WebM格式处理说明：
• WebM是现代音频格式，但百度语音识别API不直接支持
• 需要FFmpeg工具进行格式转换
• 当前系统FFmpeg未正确配置

 建议解决方案：
1. 🎯 重新录音：在录音界面选择WAV格式（推荐）
2. 🔄 在线转换：使用在线工具转换为WAV后重新上传
3. ✍️ 手动转录：直接输入音频内容

🛠️ 技术提示：
• 在线转换工具：CloudConvert、Convertio等
• 命令行：ffmpeg -i input.webm -ar 16000 -ac 1 output.wav""",

                        'm4a': f"""M4A音频文件 '{file_info}' (大小: {file_size/1024:.1f}KB) 已上传成功。

📋 M4A格式处理说明：
• M4A是常见的音频格式，百度API理论上支持
• 但需要FFmpeg工具进行预处理
• 当前系统FFmpeg未正确配置

 建议解决方案：
1. 🎯 重新录音：在录音界面系统会自动选择最佳格式
2. 🔄 格式转换：将M4A转换为WAV格式后重新上传
3. ✍️ 手动转录：直接输入音频内容

🛠️ 技术提示：
• 推荐格式：WAV (16kHz, 单声道)
• 在线转换：免费音频转换工具""",

                        'ogg': f"""OGG音频文件 '{file_info}' (大小: {file_size/1024:.1f}KB) 已上传成功。

📋 OGG格式处理说明：
• OGG是开源音频格式
• 需要FFmpeg工具进行格式转换
• 当前系统FFmpeg未正确配置

 建议解决方案：
1. 🎯 重新录音：使用WAV格式获得最佳兼容性
2. 🔄 格式转换：转换为WAV格式后重新上传
3. ✍️ 手动转录：直接输入音频内容""",

                        'flac': f"""FLAC音频文件 '{file_info}' (大小: {file_size/1024:.1f}KB) 已上传成功。

📋 FLAC格式处理说明：
• FLAC是高质量无损音频格式
• 需要FFmpeg工具进行格式转换
• 当前系统FFmpeg未正确配置

 建议解决方案：
1. 🎯 重新录音：使用WAV格式（推荐）
2. 🔄 格式转换：保持高质量转换为WAV
3. ✍️ 手动转录：直接输入音频内容"""
                    }
                    
                    message = format_messages.get(file_type.lower())
                    if message:
                        return message
                    else:
                        return f"""音频文件 '{file_info}' ({file_type.upper()}格式, {file_size/1024:.1f}KB) 已上传成功。

当前格式需要FFmpeg支持才能处理，但系统FFmpeg未正确配置。

建议：
1. 转换为WAV格式后重新上传
2. 使用录音功能直接录制WAV格式
3. 手动转录音频内容"""
                        
                except Exception as e:
                    current_app.logger.error(f"Error generating fallback message: {e}")
                    return None
            
            # 在发送到百度API之前进行最终检查
            if len(audio_data) < 1000:
                current_app.logger.error(f"Audio data too small for recognition: {len(audio_data)} bytes")
                return None
            
            if len(audio_data) > 10 * 1024 * 1024:  # 10MB
                current_app.logger.error(f"Audio data too large for recognition: {len(audio_data)} bytes")
                return None
            
            # 调用百度语音识别API
            current_app.logger.info(f"Calling Baidu ASR API with {len(audio_data)} bytes of audio data")
            result = client.asr(
                audio_data,
                'wav',    # PCM WAV格式
                16000,    # 16kHz采样率
                {
                    'dev_pid': 1537,  # 普通话(支持简单的英文识别)
                }
            )
            
            current_app.logger.info(f"Baidu API response: {result}")
            
            # 解析结果
            if result.get('err_no') == 0:
                # 识别成功
                text_result = ''.join(result.get('result', []))
                current_app.logger.info(f"Speech recognition successful: {text_result[:100]}...")
                
                # 检查结果是否为空或只有空白字符
                if text_result and text_result.strip():
                    return text_result
                else:
                    current_app.logger.warning("Speech recognition returned empty result - audio may be silent or unclear")
                    
                    # 提供有用的反馈信息
                    try:
                        file_info = os.path.basename(file_path)
                        file_size = os.path.getsize(file_path)
                        return f"""音频文件 '{file_info}' (大小: {file_size/1024:.1f}KB) 已成功上传并处理。

🔍 语音识别结果：
• 百度API响应成功，但未识别到可理解的语音内容
• 可能的原因：录音为静音、音量过小、背景噪音过大或语音不清晰

 建议改进：
1. 🎤 重新录音：确保在安静环境中清晰发音
2. 🔊 检查音量：确保录音音量适中，避免过小或过大
3. ⏱️ 录音时长：建议录音时长在3-30秒之间
4. 📝 手动输入：直接在文本框中输入要分析的内容

🛠️ 录音技巧：
• 距离麦克风15-20厘米
• 语速适中，发音清晰
• 避免背景噪音和回声"""
                    except:
                        return f"音频文件已上传，但语音识别未返回有效内容。建议重新录音或手动输入文本内容。"
            else:
                # 识别失败，提供详细的错误信息
                error_msg = result.get('err_msg', 'Unknown error')
                error_code = result.get('err_no', 'Unknown code')
                
                # 常见错误码的中文解释
                error_explanations = {
                    3300: "音频格式不正确",
                    3301: "音频数据为空",
                    3302: "音频长度过短",
                    3303: "音频长度过长",
                    3304: "客户端网络连接断开",
                    3305: "客户端网络连接超时",
                    3307: "服务端识别出错",
                    3308: "音频过长或空",
                    3309: "音频数据问题",
                    3310: "输入的音频文件过大",
                    3311: "采样率不支持",
                    3312: "音频格式不支持",
                    3313: "音频数据库连接失败",
                    3314: "音频长度不合法",
                    3315: "音频文件无法识别"
                }
                
                explanation = error_explanations.get(error_code, "未知错误")
                current_app.logger.error(f"Baidu Speech Recognition error [{error_code}]: {error_msg} ({explanation})")
                
                # 根据错误类型给出具体建议
                if error_code in [3302, 3314]:
                    current_app.logger.info("建议：录音时长至少1秒，最多60秒")
                elif error_code in [3300, 3311, 3312]:
                    current_app.logger.info("建议：使用16kHz采样率的WAV格式")
                elif error_code == 3310:
                    current_app.logger.info("建议：压缩音频文件到10MB以下")
                elif error_code == 3307:
                    current_app.logger.info("建议：检查录音质量，确保声音清晰，无背景噪音")
                    # 对于3307错误，我们可以尝试提供一个基础的音频描述
                    current_app.logger.info("尝试提供基础音频文件信息作为fallback")
                    
                    # 返回一个基础的音频文件描述而不是None
                    try:
                        file_info = os.path.basename(file_path)
                        file_size = os.path.getsize(file_path)
                        return f"音频文件 '{file_info}' (大小: {file_size/1024:.1f}KB) 已上传，但语音识别服务暂时无法处理。建议在安静环境中重新录制或手动输入内容。"
                    except:
                        return None
                
                return None
                
        finally:
            # 清理临时文件
            if converted_file_path and os.path.exists(converted_file_path):
                try:
                    os.unlink(converted_file_path)
                    current_app.logger.info(f"Cleaned up temporary file: {converted_file_path}")
                except Exception as cleanup_error:
                    current_app.logger.warning(f"Failed to cleanup temp file: {cleanup_error}")
            
    except ImportError as e:
        if 'aip' in str(e):
            current_app.logger.error("baidu-aip library not installed")
        elif 'pydub' in str(e):
            current_app.logger.error("pydub library not installed")
        else:
            current_app.logger.error(f"Import error: {e}")
        return None
    except Exception as e:
        current_app.logger.error(f"Error calling Baidu Speech API: {e}")
        return None

